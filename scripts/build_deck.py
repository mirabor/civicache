#!/usr/bin/env python3
"""
Build the CS 2640 final presentation deck.

Design principles:
- All text in real text frames (fully editable in PowerPoint).
- Custom layouts per slide; no boilerplate "Title + 3 bullets" rooms.
- Dark research aesthetic, amber accent, no clipart, no gradients.
- Big numbers and code-style mono as the visual focal points.
- Speaker notes on every slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------- design tokens ----------
BG       = RGBColor(0x0a, 0x0e, 0x1a)
PANEL    = RGBColor(0x13, 0x18, 0x26)
PANEL_2  = RGBColor(0x1c, 0x22, 0x33)
INK      = RGBColor(0xe8, 0xea, 0xed)
INK_2    = RGBColor(0x9a, 0xa3, 0xb2)
INK_3    = RGBColor(0x5d, 0x66, 0x78)
ACCENT   = RGBColor(0xf5, 0xb3, 0x42)
BLUE     = RGBColor(0x4d, 0xab, 0xf7)
GREEN    = RGBColor(0x51, 0xcf, 0x66)
RED      = RGBColor(0xff, 0x6b, 0x6b)
RULE     = RGBColor(0x2a, 0x31, 0x42)

SANS = "Helvetica Neue"
SANS_FALLBACK = "Arial"
MONO = "Menlo"

FIG = "/Users/mirayu/civicache/docs/presentation/figures"
OUT = "/Users/mirayu/civicache/docs/presentation/CS2640_Final_Deck.pptx"

# ---------- helpers ----------
def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    # remove default shadow
    spPr = shp.fill._xPr
    return shp

def add_line(slide, x1, y1, x2, y2, color=RULE, weight=Pt(1)):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = weight
    return line

def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, italic=False, color=INK, font=SANS,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             letter_space=None, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    # populate first paragraph
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if letter_space is not None:
            # XML hack for letter spacing
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(letter_space))
    return tb

def add_runs(slide, x, y, w, h, runs_spec, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs_spec: list of (text, dict_of_props) per paragraph,
       where each paragraph is itself a list of (text, props)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para_runs in enumerate(runs_spec):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for txt, props in para_runs:
            r = p.add_run()
            r.text = txt
            r.font.name = props.get("font", SANS)
            r.font.size = Pt(props.get("size", 16))
            r.font.bold = props.get("bold", False)
            r.font.italic = props.get("italic", False)
            r.font.color.rgb = props.get("color", INK)
            ls = props.get("letter_space")
            if ls is not None:
                rPr = r._r.get_or_add_rPr()
                rPr.set("spc", str(ls))
    return tb

def page_chrome(slide, idx, total, eyebrow, sub_eyebrow=None):
    # top-left eyebrow
    add_text(slide, Inches(0.5), Inches(0.32), Inches(8), Inches(0.35),
             eyebrow, size=11, bold=True, color=ACCENT, letter_space=300)
    # top-right counter
    add_text(slide, Inches(11.0), Inches(0.32), Inches(2.0), Inches(0.35),
             f"{idx:02d}  /  {total:02d}", size=11, bold=False, color=INK_3,
             align=PP_ALIGN.RIGHT, letter_space=200, font=MONO)
    # thin top rule
    add_line(slide, Inches(0.5), Inches(0.72), Inches(12.83), Inches(0.72),
             color=RULE, weight=Pt(0.75))
    # bottom footer
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             "Cache Policy Bake-Off  ·  Mira Yu  ·  CS 2640 Spring 2026",
             size=9, color=INK_3, letter_space=150)
    if sub_eyebrow:
        add_text(slide, Inches(11.0), Inches(7.05), Inches(2.0), Inches(0.3),
                 sub_eyebrow, size=9, color=INK_3, align=PP_ALIGN.RIGHT,
                 letter_space=150)

def add_speaker_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text

def add_image(slide, path, x, y, w=None, h=None):
    if w is None:
        return slide.shapes.add_picture(path, x, y, height=h)
    if h is None:
        return slide.shapes.add_picture(path, x, y, width=w)
    return slide.shapes.add_picture(path, x, y, width=w, height=h)

# ---------- build ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]
TOTAL = 12

# ===========================================================
# SLIDE 1 — COLD OPEN HOOK
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)

# top-left eyebrow only (no chrome counter on cover)
add_text(s, Inches(0.5), Inches(0.32), Inches(8), Inches(0.35),
         "CS 2640  ·  MODERN STORAGE SYSTEMS  ·  FINAL PROJECT",
         size=11, bold=True, color=ACCENT, letter_space=300)
add_text(s, Inches(11.0), Inches(0.32), Inches(2.0), Inches(0.35),
         "01  /  12", size=11, color=INK_3, align=PP_ALIGN.RIGHT,
         letter_space=200, font=MONO)
add_line(s, Inches(0.5), Inches(0.72), Inches(12.83), Inches(0.72),
         color=RULE, weight=Pt(0.75))

# Big provocative question, not a title
add_runs(s, Inches(0.5), Inches(1.05), Inches(12.5), Inches(2.4), [
    [("Same six policies. Same α sweep. Same five seeds.\n",
      {"size": 22, "color": INK_2, "letter_space": -10})],
    [("On one workload, a ", {"size": 56, "bold": True, "color": INK}),
     ("photo-finish", {"size": 56, "bold": True, "color": ACCENT}),
     (".", {"size": 56, "bold": True, "color": INK})],
    [("On the other, a ", {"size": 56, "bold": True, "color": INK}),
     ("five-point gap", {"size": 56, "bold": True, "color": ACCENT}),
     (".", {"size": 56, "bold": True, "color": INK})],
], line_spacing=1.05)

# subtitle: the question
add_text(s, Inches(0.5), Inches(3.55), Inches(12), Inches(0.5),
         "Why?  —  and what does it tell us about caching real public-records APIs?",
         size=20, color=INK_2, font=SANS)

# hero figure container (panel)
add_rect(s, Inches(0.5), Inches(4.2), Inches(12.33), Inches(2.65),
         fill=PANEL, line=RULE, line_w=Pt(0.5))
add_image(s, f"{FIG}/hero.png", Inches(0.65), Inches(4.32),
          w=Inches(12.05))

# caption
add_text(s, Inches(0.5), Inches(6.92), Inches(12.33), Inches(0.3),
         "Miss-ratio curves on Congress (left) vs. CourtListener (right). 5 seeds, ±1σ bands.",
         size=10, color=INK_3, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(7.22), Inches(8), Inches(0.3),
         "Mira Yu  ·  10-minute final talk",
         size=9, color=INK_3, letter_space=150)
add_text(s, Inches(11.0), Inches(7.22), Inches(2.0), Inches(0.3),
         "civicache",
         size=9, color=INK_3, align=PP_ALIGN.RIGHT, letter_space=200, font=MONO)

add_speaker_notes(s, """COLD OPEN — DO NOT START WITH "HI MY NAME IS".

Walk on stage and read the slide:
"Same six cache eviction policies. Same alpha sweep. Same five random seeds. Two different real workloads.

On one — a photo-finish, gaps smaller than seed noise.
On the other — a stable five-percentage-point gap at every alpha.

Why? That's the question I want to answer in the next ten minutes."

Then point at the hero figure: "Left panel — Congress.gov bills. Right panel — CourtListener court documents. Same axes, same grid. The Congress curves collapse into a narrow band. The Court curves fan into clean strata. The policies didn't change. The workloads did. The interesting question isn't WHICH policy won — it's WHY."

Pause. Then: "I'm Mira Yu. Let me show you."

(~50 sec)""")

# ===========================================================
# SLIDE 2 — THE SIX POLICIES (compact card grid)
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)
page_chrome(s, 2, TOTAL, "THE CONTESTANTS  ·  SIX EVICTION POLICIES")

add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.7),
         "Pure recency on the left  →  recency + frequency on the right.",
         size=24, bold=True, color=INK)
add_text(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.4),
         "Each policy adds exactly one mechanism over the one before it.",
         size=14, color=INK_2)

# six cards (3 columns × 2 rows)
cards = [
    ("LRU",       "Least Recently Used",     "Doubly-linked list + hash map. Promote on hit. Textbook baseline.",                        "1965",             INK_2),
    ("FIFO",      "First-In First-Out",      "No promotion. Eviction order = insertion order. Cheapest possible.",                       "Trivial",          INK_2),
    ("CLOCK",     "Clock / Second-Chance",   "Visited-bit + rotating hand. LRU approximation in O(1) amortized.",                        "Corbató '69",      INK_2),
    ("S3-FIFO",   "Three-queue FIFO",        "10% small-FIFO probation + 90% main + ghost. No LRU at all.",                              "Yang SOSP '23",    BLUE),
    ("SIEVE",     "Lazy-promotion FIFO",     "Single queue. Visited bit. Hand pointer evicts first unvisited key.",                      "Zhang NSDI '24",   GREEN),
    ("W-TinyLFU", "Frequency-aware admission", "Count-Min Sketch + admission filter + 1% window-LRU. Admits only if frequency wins.", "Einziger '17 / Caffeine", ACCENT),
]
card_w = Inches(4.05); card_h = Inches(2.45)
xs = [Inches(0.5), Inches(4.65), Inches(8.80)]
ys = [Inches(2.20), Inches(4.75)]

for i, (name, sub, body, ref, accent) in enumerate(cards):
    cx = xs[i % 3]; cy = ys[i // 3]
    add_rect(s, cx, cy, card_w, card_h, fill=PANEL, line=RULE, line_w=Pt(0.5))
    # left accent bar
    add_rect(s, cx, cy, Inches(0.06), card_h, fill=accent)
    # number
    add_text(s, cx + Inches(0.30), cy + Inches(0.18), Inches(1.0), Inches(0.4),
             f"0{i+1}", size=14, color=INK_3, font=MONO, letter_space=200)
    # name
    add_text(s, cx + Inches(0.30), cy + Inches(0.50), Inches(3.6), Inches(0.6),
             name, size=28, bold=True, color=INK)
    # sub
    add_text(s, cx + Inches(0.30), cy + Inches(1.05), Inches(3.6), Inches(0.4),
             sub, size=13, bold=True, color=accent, letter_space=120)
    # body
    add_text(s, cx + Inches(0.30), cy + Inches(1.42), Inches(3.55), Inches(0.85),
             body, size=12, color=INK_2, line_spacing=1.30)
    # ref pinned bottom
    add_text(s, cx + Inches(0.30), cy + Inches(2.10), Inches(3.55), Inches(0.30),
             ref, size=10, color=INK_3, font=MONO, letter_space=100)

add_speaker_notes(s, """THE CONTESTANTS — keep this fast, ~40 seconds.

"Six policies. I built each in C++17 from scratch — no third-party caching libraries — about 1,800 lines.

The left column is pure RECENCY. LRU you know. FIFO doesn't even promote on hit. CLOCK is LRU-approximation with a visited bit and a hand.

The right column is the recent literature. S3-FIFO from Yang SOSP'23 — three FIFO queues, no LRU at all. SIEVE from Zhang NSDI'24 — even simpler, one queue, one bit, one pointer. And W-TinyLFU — the only one that explicitly tracks FREQUENCY, with a Count-Min Sketch and an admission filter.

Reading left to right, each policy adds exactly one mechanism. That's important — it lets us isolate which mechanism is actually doing work."

Move on.""")

# ===========================================================
# SLIDE 3 — THE WORKLOAD COLLECTION PROBLEM (PROF FEEDBACK 1)
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)
page_chrome(s, 3, TOTAL, "ADDRESSING THE MIDPOINT FEEDBACK  ·  WORKLOAD COLLECTION")

add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.7),
         "How do you trace public-records APIs that publish no production logs?",
         size=24, bold=True, color=INK)

# Quote panel — professor feedback
add_rect(s, Inches(0.5), Inches(1.70), Inches(12.33), Inches(1.10),
         fill=PANEL_2, line=ACCENT, line_w=Pt(0.5))
add_text(s, Inches(0.7), Inches(1.82), Inches(2), Inches(0.30),
         "MIDPOINT FEEDBACK", size=10, bold=True, color=ACCENT, letter_space=250)
add_runs(s, Inches(0.7), Inches(2.10), Inches(11.93), Inches(0.65), [
    [("“", {"size": 24, "color": ACCENT, "italic": True}),
     ("The weighted-request idea is clever, but a real captured trace would be ideal. ",
      {"size": 16, "italic": True, "color": INK}),
     ("Find a way to get closer to a production trace.", {"size": 16, "italic": True, "bold": True, "color": INK}),
     ("”", {"size": 24, "color": ACCENT, "italic": True})],
])

# 3-column "what I tried" timeline
col_w = Inches(4.05); col_h = Inches(3.40); cy = Inches(3.10)
cols = [
    ("ATTEMPT 1",  "PACER bulk download",
     "Authoritative federal court trace. Real production traffic.",
     "Blocked: $500–$2000 paywall + per-page fees. Infeasible at 20K-request scale.",
     RED, "✗"),
    ("ATTEMPT 2",  "Synthetic Zipf only",
     "Generate request stream with α∈{0.6,…,1.2}. Trivially controlled.",
     "Doesn't answer the prof's concern. No real key space, no real sizes, no real structure.",
     INK_3, "△"),
    ("WHAT I SHIPPED", "Two real REST traces",
     "Congress.gov v3 (20.7K req) + CourtListener REST v4 (20K req). Real keys, real sizes, real timing.",
     "Plus replay-Zipf overlay so policy comparison is well-posed (next slide).",
     GREEN, "✓"),
]
for i, (tag, name, what, caveat, accent, glyph) in enumerate(cols):
    cx = Inches(0.5 + i * 4.15)
    add_rect(s, cx, cy, col_w, col_h, fill=PANEL, line=RULE, line_w=Pt(0.5))
    # glyph circle
    add_rect(s, cx + Inches(0.3), cy + Inches(0.30), Inches(0.5), Inches(0.5),
             fill=accent)
    add_text(s, cx + Inches(0.3), cy + Inches(0.31), Inches(0.5), Inches(0.5),
             glyph, size=22, bold=True, color=BG, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx + Inches(0.95), cy + Inches(0.36), Inches(3), Inches(0.3),
             tag, size=10, bold=True, color=accent, letter_space=200)
    add_text(s, cx + Inches(0.30), cy + Inches(0.95), Inches(3.6), Inches(0.55),
             name, size=22, bold=True, color=INK)
    add_text(s, cx + Inches(0.30), cy + Inches(1.55), Inches(3.6), Inches(0.95),
             what, size=12, color=INK_2, line_spacing=1.30)
    add_line(s, cx + Inches(0.30), cy + Inches(2.45),
             cx + Inches(3.75), cy + Inches(2.45), color=RULE, weight=Pt(0.5))
    add_text(s, cx + Inches(0.30), cy + Inches(2.55), Inches(3.6), Inches(0.85),
             caveat, size=11, color=INK_2, line_spacing=1.35, font=SANS)

# bottom take
add_text(s, Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.4),
         "Result: the strongest real-trace evidence the budget allowed — not a perfect production capture, but no fabrication either.",
         size=13, color=INK_2)

add_speaker_notes(s, """ADDRESSING PROF FEEDBACK — this is rehearsed.

"The professor's midpoint feedback was clear: weighted-request synthesis is clever but a real trace is better. I took that seriously.

PACER would have been ideal — federal court system, real production traffic — but the paywall is five hundred to two thousand dollars and there's no academic carve-out. Infeasible.

So I went two directions at once. Congress.gov publishes a v3 REST API with bills, amendments, and votes — twenty thousand seven hundred requests at 1.2-second pacing with backoff. CourtListener publishes a REST v4 over the same body of judicial documents PACER charges for — twenty thousand requests, ran overnight, real cluster IDs, real document sizes.

Both are REAL traces. Real keys. Real byte sizes. Real document structure. They're client-generated query streams against the public API — that's what the API exposes. There's no production user-traffic capture for either of these because no public archive exists.

I'll show next slide why I had to layer one more thing on top to make the policy comparison well-posed."

(~55 sec — important the prof hears you took the feedback seriously)""")

# ===========================================================
# SLIDE 4 — WHY REPLAY-ZIPF (the OHW problem + solution)
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)
page_chrome(s, 4, TOTAL, "BUT THE RAW TRACES HAD A PROBLEM")

add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.7),
         "Both raw traces are ~99% one-hit-wonders.",
         size=28, bold=True, color=INK)
add_text(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.5),
         "On a captured query log, every key is fetched once. There's nothing for any policy to do well — or badly — at.",
         size=15, color=INK_2)

# big-stat row
stat_w = Inches(3.95); stat_h = Inches(1.65); ny = Inches(2.40)
stats = [
    ("0.989", "Congress OHW ratio", "18,970 unique keys of 20,692 total — 91.7% of keys are seen exactly once.", INK),
    ("1.000", "Court OHW ratio (10% window)",  "15,018 unique keys of 20,000 total — every policy degenerates to FIFO on raw replay.", INK),
    ("REPLAY-ZIPF", "The fix",        "Keep real keys + real sizes. Resample popularity ranks under controlled α. Compare policies under a well-posed signal.", ACCENT),
]
for i, (n, label, body, color) in enumerate(stats):
    cx = Inches(0.5 + i * 4.10)
    add_rect(s, cx, ny, stat_w, stat_h, fill=PANEL, line=RULE, line_w=Pt(0.5))
    add_text(s, cx + Inches(0.30), ny + Inches(0.20), Inches(3.6), Inches(0.85),
             n, size=44 if n != "REPLAY-ZIPF" else 22, bold=True, color=color,
             font=MONO if n != "REPLAY-ZIPF" else SANS, letter_space=-20)
    add_text(s, cx + Inches(0.30), ny + Inches(1.00), Inches(3.6), Inches(0.30),
             label, size=11, bold=True, color=INK_2, letter_space=150)
    add_text(s, cx + Inches(0.30), ny + Inches(1.30), Inches(3.6), Inches(0.50),
             body, size=10, color=INK_2, line_spacing=1.35)

# Method box
add_rect(s, Inches(0.5), Inches(4.30), Inches(12.33), Inches(2.40),
         fill=PANEL_2, line=RULE, line_w=Pt(0.5))
add_text(s, Inches(0.7), Inches(4.45), Inches(11.93), Inches(0.4),
         "REPLAY-ZIPF METHOD  ·  WHY IT'S DEFENSIBLE", size=11, bold=True,
         color=ACCENT, letter_space=250)

add_runs(s, Inches(0.7), Inches(4.85), Inches(11.93), Inches(1.85), [
    [("Step 1.", {"size": 14, "bold": True, "color": INK}),
     ("  Capture real trace: real object keys, real byte sizes, real endpoint mix. ",
      {"size": 14, "color": INK_2}),
     ("[Congress + Court]", {"size": 14, "color": INK_3, "font": MONO})],
    [("Step 2.", {"size": 14, "bold": True, "color": INK}),
     ("  Fit MLE Zipf α to the raw popularity. Discover Congress α=0.23, Court α=1.03 — they diverge already.",
      {"size": 14, "color": INK_2})],
    [("Step 3.", {"size": 14, "bold": True, "color": INK}),
     ("  Resample popularity ranks under controlled α∈{0.6,…,1.2}. Sweep all six policies × five seeds × every α.",
      {"size": 14, "color": INK_2})],
    [("Why it works:", {"size": 14, "bold": True, "color": ACCENT}),
     ("  Differences in policy behavior trace back to ", {"size": 14, "color": INK_2}),
     ("the trace's real object population and size distribution",
      {"size": 14, "bold": True, "color": INK}),
     (" — not to the request-interval pattern, which the public APIs don't expose anyway.",
      {"size": 14, "color": INK_2})],
])

add_speaker_notes(s, """THE OHW PROBLEM AND WHY REPLAY-ZIPF.

"Here's the thing nobody told me until I had the data: both raw traces are dominated by one-hit wonders. Congress is 99 percent. Court is 100 percent in a 10-percent rolling window.

If I just replay those traces, every admission becomes a miss, every eviction is a key that would never be re-accessed anyway. All six policies converge to FIFO. There's nothing to measure.

So I do replay-Zipf. Step one: capture the real trace — real keys, real byte sizes. Step two: measure the raw alpha. Step three: resample the popularity ranks under controlled alpha values from 0.6 to 1.2.

The KEYS are real. The SIZES are real. The popularity DISTRIBUTION is the controlled variable.

Why is this defensible? Because the differences I'm about to show you trace back to the trace's real object population and real size distribution — not to the request-interval pattern, which the public APIs don't expose anyway. There's no production user log to compare against."

(~50 sec)""")

# ===========================================================
# SLIDE 5 — THE TWO WORKLOADS DIVERGE (data + figures)
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)
page_chrome(s, 5, TOTAL, "THE TWO WORKLOADS DIVERGE  ·  GROUNDED IN RAW STATS")

add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.7),
         "Two real APIs.  One looks uniform.  One has a 462 KB outlier.",
         size=24, bold=True, color=INK)

# split table
add_rect(s, Inches(0.5), Inches(1.85), Inches(6.40), Inches(4.95),
         fill=PANEL, line=RULE, line_w=Pt(0.5))

add_text(s, Inches(0.7), Inches(2.00), Inches(6), Inches(0.3),
         "RAW TRACE STATS", size=11, bold=True, color=ACCENT, letter_space=250)

# table headers
hdr_y = Inches(2.40)
add_text(s, Inches(0.7),  hdr_y, Inches(2.5), Inches(0.3), "Metric",
         size=11, bold=True, color=INK_2, letter_space=120)
add_text(s, Inches(3.4),  hdr_y, Inches(1.6), Inches(0.3), "Congress",
         size=11, bold=True, color=BLUE, letter_space=120, align=PP_ALIGN.RIGHT)
add_text(s, Inches(5.05), hdr_y, Inches(1.7), Inches(0.3), "Court",
         size=11, bold=True, color=ACCENT, letter_space=120, align=PP_ALIGN.RIGHT)
add_line(s, Inches(0.7), Inches(2.72), Inches(6.75), Inches(2.72),
         color=RULE, weight=Pt(0.75))

rows = [
    ("Total requests",        "20,692",      "20,000"),
    ("Unique objects",        "18,970",      "15,018"),
    ("Zipf α (MLE)",          "0.231",       "1.028"),
    ("OHW ratio",             "0.989",       "1.000"),
    ("Median size (bytes)",   "231",         "1,381"),
    ("Max size (bytes)",      "6,700",       "462,490"),
    ("Max / median ratio",    "29×",         "335×"),
    ("Unique-key fraction",   "91.7%",       "75.1%"),
]
for i, (m, c, ct) in enumerate(rows):
    ry = Inches(2.85 + i * 0.46)
    add_text(s, Inches(0.7),  ry, Inches(2.5), Inches(0.40), m,
             size=12, color=INK)
    add_text(s, Inches(3.4),  ry, Inches(1.6), Inches(0.40), c,
             size=12, color=INK, font=MONO, align=PP_ALIGN.RIGHT)
    # highlight the divergent rows (alpha + max size)
    is_div = m in ("Zipf α (MLE)", "Max size (bytes)", "Max / median ratio")
    add_text(s, Inches(5.05), ry, Inches(1.7), Inches(0.40), ct,
             size=12, color=ACCENT if is_div else INK, font=MONO,
             bold=is_div, align=PP_ALIGN.RIGHT)

# Right-side figures
add_rect(s, Inches(7.10), Inches(1.85), Inches(5.73), Inches(2.30),
         fill=PANEL, line=RULE, line_w=Pt(0.5))
add_text(s, Inches(7.30), Inches(1.95), Inches(5), Inches(0.3),
         "CONGRESS  ·  α_raw = 0.23  ·  ESSENTIALLY UNIFORM",
         size=10, bold=True, color=BLUE, letter_space=180)
add_image(s, f"{FIG}/workload_congress.png", Inches(7.25), Inches(2.30),
          w=Inches(5.45))

add_rect(s, Inches(7.10), Inches(4.30), Inches(5.73), Inches(2.50),
         fill=PANEL, line=RULE, line_w=Pt(0.5))
add_text(s, Inches(7.30), Inches(4.40), Inches(5), Inches(0.3),
         "COURT  ·  α_raw = 1.03  ·  GENUINE LONG TAIL",
         size=10, bold=True, color=ACCENT, letter_space=180)
add_image(s, f"{FIG}/workload_court.png", Inches(7.25), Inches(4.75),
          w=Inches(5.45))

add_speaker_notes(s, """THE WORKLOADS DIVERGE — keep ~50 sec.

"The data tells the story.

Congress.gov: alpha 0.23. Median object 231 bytes. Max object 6.7 KB. Almost 92 percent of keys appear exactly once. Effectively uniform random over the bill / amendment / vote key space.

CourtListener: alpha 1.03 — that's a 4.4× higher skew, MLE-fit, the same way the Clauset 2009 power-law paper does it. Median 1.4 KB, but the MAX is 462 KILOBYTES. That's 335 times the median. One single document — a long opinion that's full text in JSON — dominates the byte-distribution.

Why does Court have real locality and Congress doesn't? Because CourtListener serves judicial documents that get cited across briefs, dockets that get refreshed during active cases, cluster IDs that recur across related searches. Congress.gov has no front-end aggregator that would re-request the same bill; nobody fetches bill 1234 immediately after bill 5678.

These two raw traces are the foundation for everything that follows."

(~55 sec)""")

# ===========================================================
# SLIDE 6 — HEADLINE RESULT + THE PUZZLE
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)
page_chrome(s, 6, TOTAL, "HEADLINE RESULT  ·  AND THE PUZZLE THAT REMAINS")

# Big number left
add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.6),
         "W-TinyLFU wins 11 of 12 cells.  But not on the cell you'd expect.",
         size=24, bold=True, color=INK)

# Cell-grid: 6 cache fractions × 2 workloads = 12
add_rect(s, Inches(0.5), Inches(1.80), Inches(6.40), Inches(4.40),
         fill=PANEL, line=RULE, line_w=Pt(0.5))
add_text(s, Inches(0.7), Inches(1.93), Inches(6), Inches(0.3),
         "WIN GRID  ·  6 CACHE FRACTIONS × 2 WORKLOADS",
         size=11, bold=True, color=ACCENT, letter_space=200)

# headers: cache fractions
fracs = ["0.001", "0.005", "0.010", "0.020", "0.050", "0.100"]
cell_x0 = Inches(2.10); cell_y0 = Inches(2.90)
cell_w = Inches(0.75); cell_h = Inches(1.10); gap = Inches(0.04)

# row labels
add_text(s, Inches(0.7), Inches(2.50), Inches(1.4), Inches(0.30), "cache_frac:",
         size=10, color=INK_3, font=MONO, letter_space=120)
for i, f in enumerate(fracs):
    add_text(s, cell_x0 + (cell_w + gap) * i, Inches(2.55), cell_w, Inches(0.30),
             f, size=11, color=INK_2, font=MONO, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), cell_y0 + Inches(0.40), Inches(1.4), Inches(0.30),
         "Congress", size=12, bold=True, color=BLUE)
add_text(s, Inches(0.7), cell_y0 + cell_h + gap + Inches(0.40),
         Inches(1.4), Inches(0.30),
         "Court", size=12, bold=True, color=ACCENT)

# winners per cell — Congress: WTLFU,WTLFU,WTLFU,SIEVE,WTLFU,WTLFU; Court: all WTLFU
congress_winners = ["W-TLFU", "W-TLFU", "W-TLFU", "SIEVE", "W-TLFU", "W-TLFU"]
court_winners    = ["W-TLFU"] * 6
for i, w in enumerate(congress_winners):
    cx = cell_x0 + (cell_w + gap) * i
    color = ACCENT if w == "W-TLFU" else GREEN
    add_rect(s, cx, cell_y0, cell_w, cell_h, fill=color)
    add_text(s, cx, cell_y0 + Inches(0.30), cell_w, Inches(0.5),
             w, size=10, bold=True, color=BG, font=MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, w in enumerate(court_winners):
    cx = cell_x0 + (cell_w + gap) * i
    color = ACCENT
    add_rect(s, cx, cell_y0 + cell_h + gap, cell_w, cell_h, fill=color)
    add_text(s, cx, cell_y0 + cell_h + gap + Inches(0.30), cell_w, Inches(0.5),
             "W-TLFU", size=10, bold=True, color=BG, font=MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# legend
add_rect(s, Inches(0.7), Inches(5.40), Inches(0.18), Inches(0.18), fill=ACCENT)
add_text(s, Inches(0.95), Inches(5.35), Inches(2), Inches(0.30),
         "W-TinyLFU wins", size=11, color=INK_2)
add_rect(s, Inches(2.7), Inches(5.40), Inches(0.18), Inches(0.18), fill=GREEN)
add_text(s, Inches(2.95), Inches(5.35), Inches(2), Inches(0.30),
         "SIEVE wins", size=11, color=INK_2)

add_text(s, Inches(0.7), Inches(5.85), Inches(6), Inches(0.30),
         "Welch's t-test, p < 0.05, 5 seeds, ±1σ bands.",
         size=10, color=INK_3, font=MONO)

# Right side: the puzzle
add_rect(s, Inches(7.10), Inches(1.80), Inches(5.73), Inches(4.40),
         fill=PANEL_2, line=ACCENT, line_w=Pt(0.5))
add_text(s, Inches(7.30), Inches(1.93), Inches(5), Inches(0.3),
         "THE PUZZLE", size=11, bold=True, color=ACCENT, letter_space=250)

add_runs(s, Inches(7.30), Inches(2.30), Inches(5.40), Inches(3.85), [
    [("11 / 12", {"size": 64, "bold": True, "color": INK, "font": MONO,
                  "letter_space": -40})],
    [("cells where W-TinyLFU dominates",
      {"size": 13, "color": INK_2})],
    [(" ", {"size": 8})],
    [("but on Congress, ", {"size": 16, "color": INK_2}),
     ("the W-TinyLFU vs SIEVE gap is always within ±0.005 — less than one σ.",
      {"size": 16, "color": INK})],
    [(" ", {"size": 8})],
    [("On Court, the same two policies separate by ",
      {"size": 16, "color": INK_2}),
     ("a stable 4–5 pp at every α.",
      {"size": 16, "bold": True, "color": ACCENT})],
    [(" ", {"size": 10})],
    [("Why does the same admission filter dominate one workload and tie on the other?",
      {"size": 15, "italic": True, "color": ACCENT})],
])

add_speaker_notes(s, """HEADLINE — set up the puzzle.

"The headline number — W-TinyLFU wins 11 of 12 miss-ratio cells. Six cache fractions, two workloads. Welch's t-test, p less than 0.05, five seeds.

But look at the grid: the one cell W-TinyLFU LOSES is on Congress at cache_frac 0.020. SIEVE — that's the simpler policy, the visited-bit FIFO from NSDI'24 — beats it.

And it's not just that one cell. Across the entire alpha sweep on Congress, the W-TinyLFU vs SIEVE gap is always within plus or minus 0.005. Less than one seed-sigma. Statistically tied.

On Court — the SAME two policies, the SAME implementation — separate by a stable four to five percentage points at every alpha.

So here's the puzzle: how can the same admission filter dominate one workload and tie on the other? The TinyLFU paper would predict it dominates everywhere skew is high. It doesn't.

Next slide is the most important slide in the talk."

(~50 sec — build tension here, the prof asked specifically for the WHY)""")

# ===========================================================
# SLIDE 7 — MECHANISM CENTERPIECE (the WHY)
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)
page_chrome(s, 7, TOTAL, "MECHANISM  ·  WHY THE GAP DEPENDS ON THE WORKLOAD")

add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.7),
         "Admission-filter wins are gated by raw-trace frequency gradient.",
         size=26, bold=True, color=INK)
add_text(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.4),
         "SIEVE's visited-bit is functionally a poor-man's admission filter.  CMS only beats it when there's a gradient to capture.",
         size=14, color=INK_2)

# left: Congress card
add_rect(s, Inches(0.5), Inches(2.20), Inches(6.10), Inches(4.65),
         fill=PANEL, line=BLUE, line_w=Pt(0.75))
add_text(s, Inches(0.70), Inches(2.35), Inches(6), Inches(0.30),
         "CONGRESS  ·  α_raw = 0.23",
         size=11, bold=True, color=BLUE, letter_space=200)
add_text(s, Inches(0.70), Inches(2.65), Inches(5.7), Inches(0.50),
         "Both signals collapse to binary.",
         size=22, bold=True, color=INK)
# what each policy "sees"
add_text(s, Inches(0.70), Inches(3.30), Inches(0.7), Inches(0.30),
         "CMS:", size=11, bold=True, color=INK_2, font=MONO, letter_space=150)
add_text(s, Inches(1.50), Inches(3.30), Inches(4.8), Inches(0.30),
         "[1, 1, 1, 1, 1, 2, 1, 1, 1] →  near-uniform counters",
         size=12, color=INK_3, font=MONO)

add_text(s, Inches(0.70), Inches(3.65), Inches(0.9), Inches(0.30),
         "SIEVE:", size=11, bold=True, color=INK_2, font=MONO, letter_space=150)
add_text(s, Inches(1.65), Inches(3.65), Inches(4.7), Inches(0.30),
         "visited-bit ∈ {0, 1}  →  also binary",
         size=12, color=INK_3, font=MONO)

add_line(s, Inches(0.70), Inches(4.10), Inches(6.40), Inches(4.10),
         color=RULE, weight=Pt(0.5))

add_runs(s, Inches(0.70), Inches(4.25), Inches(5.7), Inches(2.50), [
    [("→", {"size": 22, "color": ACCENT, "bold": True}),
     ("  Both policies converge to the same predicate:",
      {"size": 14, "color": INK_2})],
    [("    “has this key been seen more than once?”",
      {"size": 14, "italic": True, "color": INK})],
    [(" ", {"size": 10})],
    [("Result:  ", {"size": 14, "bold": True, "color": INK})],
    [("Congress gap (SIEVE − W-TLFU) ∈ [−0.005, +0.002]  across α∈[0.6,1.5]",
      {"size": 12, "color": INK_2, "font": MONO})],
    [("Statistically tied — gap < seed σ at every point.",
      {"size": 13, "color": ACCENT, "bold": True})],
])

# right: Court card
add_rect(s, Inches(6.73), Inches(2.20), Inches(6.10), Inches(4.65),
         fill=PANEL, line=ACCENT, line_w=Pt(0.75))
add_text(s, Inches(6.93), Inches(2.35), Inches(6), Inches(0.30),
         "COURT  ·  α_raw = 1.03  +  462 KB OUTLIER",
         size=11, bold=True, color=ACCENT, letter_space=200)
add_text(s, Inches(6.93), Inches(2.65), Inches(5.7), Inches(0.50),
         "CMS captures a real gradient.",
         size=22, bold=True, color=INK)

add_text(s, Inches(6.93), Inches(3.30), Inches(0.7), Inches(0.30),
         "CMS:", size=11, bold=True, color=INK_2, font=MONO, letter_space=150)
add_text(s, Inches(7.73), Inches(3.30), Inches(4.8), Inches(0.30),
         "[47, 12, 8, 5, 2, 2, 1, 1, …] →  rich frequency strata",
         size=12, color=ACCENT, font=MONO, bold=True)

add_text(s, Inches(6.93), Inches(3.65), Inches(0.9), Inches(0.30),
         "SIEVE:", size=11, bold=True, color=INK_2, font=MONO, letter_space=150)
add_text(s, Inches(7.88), Inches(3.65), Inches(4.7), Inches(0.30),
         "visited-bit ∈ {0, 1}  →  loses the gradient",
         size=12, color=INK_3, font=MONO)

add_line(s, Inches(6.93), Inches(4.10), Inches(12.63), Inches(4.10),
         color=RULE, weight=Pt(0.5))

add_runs(s, Inches(6.93), Inches(4.25), Inches(5.7), Inches(2.50), [
    [("→", {"size": 22, "color": ACCENT, "bold": True}),
     ("  Admission filter rejects the 462 KB outlier on first sight.",
      {"size": 14, "color": INK_2})],
    [("    SIEVE admits it, evicts ~5–10% of capacity in one shot.",
      {"size": 13, "italic": True, "color": INK_2})],
    [(" ", {"size": 10})],
    [("Result:  ", {"size": 14, "bold": True, "color": INK})],
    [("Court gap (SIEVE − W-TLFU) ∈ [+0.036, +0.054]  across α∈[0.6,1.2]",
      {"size": 12, "color": INK_2, "font": MONO})],
    [("Stable +4 to +5 pp gap, every α, every seed.",
      {"size": 13, "color": ACCENT, "bold": True})],
])

add_speaker_notes(s, """THE CENTERPIECE — slow down. This is the prof's #1 ask.

"This is the answer.

SIEVE's visited-bit is FUNCTIONALLY a poor-man's admission filter. On a miss, a key enters with visited=false. On a hit, the bit flips. The eviction scanner clears trues and evicts the first false it finds. So SIEVE evicts one-hit-wonders aggressively — same goal as TinyLFU's admission filter, just expressed in one bit instead of a Count-Min Sketch.

The DIFFERENCE between the bit and the sketch is that the BIT is binary and the SKETCH is a gradient.

On Congress — alpha 0.23, near-uniform — the CMS sees [1,1,1,1,2,1,1,1...] — almost-all-ones counters. There's no gradient to capture. The CMS collapses to the same binary predicate the visited-bit already implements: 'have I seen this key more than once?' That's why they tie.

On Court — alpha 1.03, plus that 462 KB outlier — the CMS sees [47, 12, 8, 5, 2, 1, 1...]. A real frequency stratification. The admission filter can REFUSE the 462 KB document on first sight. SIEVE has to admit it, then evict it after one cycle — and during that cycle the document is taking up 5 to 10 percent of the small cache.

That's why the same two policies tie on one workload and separate by 5pp on the other. The mechanism is the SAME. The information available to it is different.

This is the finding. The TinyLFU paper says 'W-TinyLFU wins on Zipf' — implicitly assuming the Zipf overlay is the whole story. Our data shows the RAW trace structure matters independently of the alpha overlay."

(~70 sec — most important slide, take the time)""")

# ===========================================================
# SLIDE 8 — ABLATIONS (rigor)
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)
page_chrome(s, 8, TOTAL, "ABLATIONS  ·  ISOLATING WHICH MECHANISMS MATTER")

add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.7),
         "Three controlled tests. Hold the policy fixed; vary one design choice.",
         size=24, bold=True, color=INK)

# 3-column ablation grid
ax = [Inches(0.5), Inches(4.65), Inches(8.80)]
ay = Inches(1.85); aw = Inches(4.05); ah = Inches(5.00)

ablations = [
    ("S3-FIFO  ·  small-queue ratio",
     "5%  vs  10% (paper default)  vs  20%",
     "−6.3 pp",
     "Court @ α=1.2:  5% beats 20%",
     "Recommendation:  reduce small-queue toward 5% on heavy-tailed workloads — overrides the SOSP'23 paper default.",
     ACCENT, BLUE),
    ("SIEVE  ·  lazy promotion",
     "SIEVE  vs  SIEVE-NoProm",
     "+15.4 pp",
     "Congress @ α=1.0:  largest delta in study",
     "Confirms NSDI'24 claim:  the visited-bit is what makes SIEVE not-FIFO.  Load-bearing.",
     ACCENT, GREEN),
    ("W-TinyLFU  ·  Doorkeeper",
     "W-TLFU  vs  W-TLFU + Doorkeeper",
     "± 0.7 pp",
     "Direction-swaps at α extremes",
     "Confirms Caffeine production decision:  Doorkeeper is a marginal hedge, not a default.",
     INK_2, INK_3),
]
for i, (title, sweep, delta, where, takeaway, dcolor, accent) in enumerate(ablations):
    cx = ax[i]
    add_rect(s, cx, ay, aw, ah, fill=PANEL, line=RULE, line_w=Pt(0.5))
    add_rect(s, cx, ay, Inches(0.06), ah, fill=accent)
    add_text(s, cx + Inches(0.30), ay + Inches(0.20), Inches(3.6), Inches(0.40),
             f"ABLATION 0{i+1}", size=10, bold=True, color=accent, letter_space=200)
    add_text(s, cx + Inches(0.30), ay + Inches(0.55), Inches(3.6), Inches(0.85),
             title, size=18, bold=True, color=INK)
    add_text(s, cx + Inches(0.30), ay + Inches(1.40), Inches(3.6), Inches(0.40),
             sweep, size=12, color=INK_2, font=MONO)
    add_line(s, cx + Inches(0.30), ay + Inches(1.85),
             cx + Inches(3.75), ay + Inches(1.85), color=RULE, weight=Pt(0.5))

    add_text(s, cx + Inches(0.30), ay + Inches(2.00), Inches(3.6), Inches(0.30),
             "PEAK  Δ", size=10, bold=True, color=INK_3, letter_space=200)
    add_text(s, cx + Inches(0.30), ay + Inches(2.30), Inches(3.6), Inches(1.00),
             delta, size=42, bold=True, color=dcolor, font=MONO,
             letter_space=-30)
    add_text(s, cx + Inches(0.30), ay + Inches(3.30), Inches(3.6), Inches(0.40),
             where, size=11, color=INK_2, italic=False)
    add_line(s, cx + Inches(0.30), ay + Inches(3.75),
             cx + Inches(3.75), ay + Inches(3.75), color=RULE, weight=Pt(0.5))
    add_text(s, cx + Inches(0.30), ay + Inches(3.85), Inches(3.6), Inches(1.10),
             takeaway, size=11, color=INK, line_spacing=1.30)

# bottom synthesis
add_rect(s, Inches(0.5), Inches(7.00), Inches(12.33), Inches(0.30),
         fill=PANEL_2, line=None)
add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.30),
         "Reading the three together: policy ordering is a property of the policy families, not of the parameter knobs we shipped.",
         size=12, color=INK_2, italic=False)

add_speaker_notes(s, """ABLATIONS — fast, ~50 sec total. Land the synthesis.

"Three controlled ablations. Hold the policy fixed, vary one design choice.

ONE — S3-FIFO small-queue ratio. The SOSP'23 paper recommends 10%. I tested 5, 10, and 20 percent. 5% beats 20% by 6.3 percentage points on Court at high alpha, monotonically. Actionable advice — reduce the small-queue ratio on heavy-tailed workloads.

TWO — SIEVE's visited-bit. I toggle promote-on-hit off. SIEVE-NoProm loses by FIFTEEN POINT FOUR percentage points on Congress at alpha 1.0. That's the largest ablation delta in the entire study. Confirms the NSDI'24 paper's claim — the visited-bit is what makes SIEVE not-FIFO.

THREE — W-TinyLFU plus or minus the Doorkeeper Bloom filter. The original TinyLFU paper proposes it; Caffeine production omits it. My data: noise. Plus or minus 0.7pp, direction-swaps at alpha extremes. Confirms Caffeine's decision.

The pattern: policy ORDERING is a property of the policy FAMILY, not of the parameter knobs. That's why I can claim the result on the previous slide is robust."

(~50 sec)""")

# ===========================================================
# SLIDE 9 — SHARDS VALIDATION (depth signal)
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)
page_chrome(s, 9, TOTAL, "SHARDS  ·  TRUSTING THE MISS-RATIO CURVES")

add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.7),
         "No exact MRC oracle exists at 1M-access scale.  How do you trust the curve?",
         size=24, bold=True, color=INK)

# Left: explanation
add_rect(s, Inches(0.5), Inches(1.85), Inches(6.20), Inches(5.00),
         fill=PANEL, line=RULE, line_w=Pt(0.5))
add_text(s, Inches(0.7), Inches(2.00), Inches(6), Inches(0.3),
         "VALIDATION VIA SELF-CONVERGENCE  ·  WALDSPURGER FAST'15",
         size=11, bold=True, color=ACCENT, letter_space=180)

add_runs(s, Inches(0.7), Inches(2.45), Inches(5.85), Inches(4.30), [
    [("Problem:", {"size": 14, "bold": True, "color": INK}),
     ("  Exact stack-distance is O(n × |unique|).  Hours on a 1M trace.",
      {"size": 14, "color": INK_2})],
    [(" ", {"size": 8})],
    [("Method:", {"size": 14, "bold": True, "color": INK}),
     ("  Hash-based sampling at 4 rates {0.01, 0.1, 1, 10}%.",
      {"size": 14, "color": INK_2})],
    [(" ", {"size": 6})],
    [("       Compare each non-reference rate vs 10% reference.",
      {"size": 13, "color": INK_2})],
    [("       If the curves agree, the sampled MRC is reliable.",
      {"size": 13, "color": INK_2})],
    [(" ", {"size": 8})],
    [("Sanity gate:", {"size": 14, "bold": True, "color": INK}),
     ("  MAE < 0.05 vs reference.",
      {"size": 14, "color": INK_2})],
    [(" ", {"size": 12})],
    [(" 1% vs 10%   →   MAE = 0.0378", {"size": 16, "color": GREEN, "font": MONO, "bold": True}),
     ("   PASS", {"size": 14, "color": GREEN, "bold": True})],
    [(" 0.1% vs 10% →   MAE = 0.0496", {"size": 14, "color": INK_2, "font": MONO}),
     ("   borderline", {"size": 12, "color": INK_2})],
    [(" 0.01%       →   81 effective samples", {"size": 14, "color": RED, "font": MONO}),
     ("   below 200 floor — flagged", {"size": 12, "color": RED})],
    [(" ", {"size": 12})],
    [("Cross-check:", {"size": 14, "bold": True, "color": INK}),
     ("  50 K-access oracle regime with exact stack distances reproduces SHARDS within tolerance — the convergence is not a 1M-only artifact.",
      {"size": 13, "color": INK_2})],
])

# Right: figure
add_rect(s, Inches(6.83), Inches(1.85), Inches(6.00), Inches(5.00),
         fill=PANEL, line=RULE, line_w=Pt(0.5))
add_text(s, Inches(7.00), Inches(1.95), Inches(5.5), Inches(0.3),
         "FOUR-RATE OVERLAY  ·  1M ZIPF α=0.8",
         size=11, bold=True, color=ACCENT, letter_space=180)
add_image(s, f"{FIG}/shards_overlay.png", Inches(6.95), Inches(2.40),
          w=Inches(5.75))

add_speaker_notes(s, """SHARDS — depth signal, ~45 sec.

"Quick rigor checkpoint. Every miss-ratio curve I just showed depends on stack-distance computation. At 1 million accesses the exact algorithm is O(n × unique keys) — hours per run. Infeasible.

Waldspurger's FAST'15 SHARDS paper introduces hash-based sampling. I run SHARDS at four sampling rates — point-zero-one, point-one, one, and ten percent — and check that the MRCs CONVERGE. If they agree, the sampled MRC is reliable.

The gate is mean-absolute-error less than 0.05 against the 10% reference. The 1%-vs-10% MAE is 0.0378 — passes. The 0.01% rate has only 81 effective samples — below my 200-sample floor — flagged but reported.

I also run a 50K-access regime where exact stack distances ARE tractable, and SHARDS reproduces the exact curve there. So the convergence isn't a 1M-only numerical artifact — it extrapolates from a workload where I can directly verify it. That's what lets me trust the 1M numbers."

(~45 sec)""")

# ===========================================================
# SLIDE 10 — BYTE-MRC TRAP
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)
page_chrome(s, 10, TOTAL, "METHODOLOGICAL WARNING  ·  BYTE-MRC ON HEAVY TAILS")

add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.7),
         "On Court, single-seed byte-miss could flip the policy ordering.",
         size=24, bold=True, color=INK)

# top stat strip
add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(1.30),
         fill=PANEL_2, line=RULE, line_w=Pt(0.5))
add_text(s, Inches(0.7), Inches(2.00), Inches(12), Inches(0.3),
         "WHY ONE OBJECT BREAKS THE METRIC", size=11, bold=True,
         color=ACCENT, letter_space=200)
add_runs(s, Inches(0.7), Inches(2.40), Inches(12), Inches(0.7), [
    [("462 KB", {"size": 28, "bold": True, "color": ACCENT, "font": MONO}),
     ("   one document.   ", {"size": 16, "color": INK_2}),
     ("0.78%", {"size": 22, "bold": True, "color": INK, "font": MONO}),
     ("  of the working set by bytes  ·  ", {"size": 16, "color": INK_2}),
     ("0.007%", {"size": 22, "bold": True, "color": INK, "font": MONO}),
     ("  of the objects by count.", {"size": 16, "color": INK_2})],
])

# table — object miss vs byte miss
add_rect(s, Inches(0.5), Inches(3.30), Inches(7.50), Inches(3.55),
         fill=PANEL, line=RULE, line_w=Pt(0.5))
add_text(s, Inches(0.7), Inches(3.45), Inches(7), Inches(0.3),
         "COURT @ cache_frac = 0.01  ·  5-SEED MEAN ± σ",
         size=11, bold=True, color=ACCENT, letter_space=180)

# headers
y = Inches(3.95)
add_text(s, Inches(0.7),  y, Inches(2.4), Inches(0.3), "Policy",
         size=11, bold=True, color=INK_2, letter_space=120)
add_text(s, Inches(3.10), y, Inches(2.3), Inches(0.3), "Object miss-ratio",
         size=11, bold=True, color=INK_2, letter_space=120)
add_text(s, Inches(5.50), y, Inches(2.3), Inches(0.3), "Byte miss-ratio",
         size=11, bold=True, color=INK_2, letter_space=120)
add_line(s, Inches(0.7), Inches(4.27), Inches(7.85), Inches(4.27),
         color=RULE, weight=Pt(0.75))

bm_rows = [
    ("W-TinyLFU", "0.728  ±  0.045", "0.611  ±  0.207", ACCENT),
    ("SIEVE",     "0.783  ±  0.045", "0.656  ±  0.205", GREEN),
    ("LRU",       "0.847  ±  0.026", "0.755  ±  0.171", INK_2),
]
for i, (p, om, bm, c) in enumerate(bm_rows):
    ry = Inches(4.42 + i * 0.55)
    add_text(s, Inches(0.7),  ry, Inches(2.4), Inches(0.4), p,
             size=14, bold=True, color=c)
    add_text(s, Inches(3.10), ry, Inches(2.3), Inches(0.4), om,
             size=13, color=INK, font=MONO)
    add_text(s, Inches(5.50), ry, Inches(2.3), Inches(0.4), bm,
             size=13, color=ACCENT if i == 0 else INK, font=MONO,
             bold=(i == 0))

add_text(s, Inches(0.7), Inches(6.20), Inches(7.5), Inches(0.6),
         "Object σ ≈ 0.03–0.05 (normal). Byte σ ≈ 0.17–0.21 (4–7× larger).",
         size=12, color=INK_2)

# Right callout
add_rect(s, Inches(8.20), Inches(3.30), Inches(4.63), Inches(3.55),
         fill=PANEL, line=RED, line_w=Pt(0.75))
add_text(s, Inches(8.40), Inches(3.45), Inches(4.5), Inches(0.3),
         "THE TRAP",  size=11, bold=True, color=RED, letter_space=250)
add_runs(s, Inches(8.40), Inches(3.85), Inches(4.30), Inches(2.90), [
    [("14.4 pp", {"size": 38, "bold": True, "color": INK, "font": MONO})],
    [("inter-policy gap (W-TLFU vs LRU)",
      {"size": 12, "color": INK_2})],
    [(" ", {"size": 10})],
    [("vs", {"size": 16, "italic": True, "color": INK_3})],
    [(" ", {"size": 10})],
    [("20+ pp", {"size": 38, "bold": True, "color": RED, "font": MONO})],
    [("seed-to-seed σ on byte-miss",
      {"size": 12, "color": INK_2})],
    [(" ", {"size": 10})],
    [("→ A single seed could report W-TinyLFU LOSING to LRU on bytes.",
      {"size": 12, "italic": True, "color": ACCENT, "bold": True})],
])

add_speaker_notes(s, """BYTE-MRC TRAP — ~40 sec, depth signal.

"One more methodological point that's worth landing.

Byte miss-ratio on Court is dominated by that single 462 KB document. It's 0.78 percent of the working set by bytes but only 0.007 percent of objects by count.

Look at the table. Object-miss sigma is 0.03 to 0.05 — normal. Byte-miss sigma is 0.17 to 0.21 — four to seven times larger.

The gap between W-TinyLFU and LRU on byte-miss is 14.4 pp. The seed-to-seed sigma is 20-plus pp. A single-seed byte-miss benchmark could easily report W-TinyLFU LOSING to LRU on bytes — purely from which Zipf rank that one outlier document happened to land at on that seed.

This isn't a problem with my measurement — it's a property of byte-miss-ratio on heavy-tailed workloads. The methodological lesson: report 5-seed CI bands on byte-MRC, or you're going to get burned. Cite this if you ever see a single-seed byte-miss benchmark that disagrees with the object-miss conclusion."

(~40 sec)""")

# ===========================================================
# SLIDE 11 — DECISION TREE (practitioner takeaway)
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)
page_chrome(s, 11, TOTAL, "PRACTITIONER DECISION TREE")

add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.7),
         "Pick the eviction policy by reading the trace, not the paper.",
         size=24, bold=True, color=INK)
add_text(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.4),
         "Run α-MLE on the raw popularity. Check the max/median size ratio. Then choose.",
         size=14, color=INK_2)

# decision tree as 4 branches
branches = [
    ("α_raw  <  0.5    AND    sizes ≈ uniform",
     "Congress-like",
     "→  SIEVE",
     "Statistically tied with W-TinyLFU across α∈[0.6, 1.5]. Prefer simpler policy: no CMS, no admission filter, just one bit + a hand pointer.",
     GREEN),
    ("α_raw  >  0.8    AND    long-tail sizes",
     "Court-like",
     "→  W-TinyLFU",
     "Stable +5 pp at every α. CMS gradient extracts information SIEVE's binary bit cannot. CMS + admission filter + 1% window-LRU is worth the implementation cost.",
     ACCENT),
    ("Catastrophic size outlier present",
     "max / median  >  100×",
     "→  W-TinyLFU + report byte-MRC with 5-seed CI",
     "Single-seed byte-MRC is dominated by which Zipf rank the outlier lands at. Always show ±σ bands.",
     RED),
    ("cache_frac  <  0.01    (small cache)",
     "any workload",
     "→  W-TinyLFU dominates everywhere",
     "+9 to +13 pp over LRU on both workloads. Admission quality dominates when one wrong admit evicts 5–10% of capacity.",
     BLUE),
]

bx = Inches(0.5); by0 = Inches(2.20); bw = Inches(12.33); bh = Inches(1.10); bgap = Inches(0.10)

for i, (cond, label, pick, expl, color) in enumerate(branches):
    by = by0 + (bh + bgap) * i
    add_rect(s, bx, by, bw, bh, fill=PANEL, line=RULE, line_w=Pt(0.5))
    add_rect(s, bx, by, Inches(0.10), bh, fill=color)

    # left zone — condition
    add_text(s, bx + Inches(0.35), by + Inches(0.10), Inches(4.5), Inches(0.40),
             cond, size=14, bold=True, color=INK, font=MONO)
    add_text(s, bx + Inches(0.35), by + Inches(0.55), Inches(4.5), Inches(0.40),
             label, size=11, color=INK_3, italic=False, letter_space=200)

    # mid — pick
    add_text(s, bx + Inches(5.10), by + Inches(0.20), Inches(3.0), Inches(0.65),
             pick, size=22, bold=True, color=color, font=MONO,
             anchor=MSO_ANCHOR.MIDDLE)

    # right — explanation
    add_text(s, bx + Inches(8.30), by + Inches(0.10), Inches(3.85), Inches(0.95),
             expl, size=11, color=INK_2, line_spacing=1.30)

# bottom rail
add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.30),
         "Workload-first thinking.  The policy choice is ~5 lines of code; the rationale is the data on slide 5.",
         size=12, italic=False, color=INK_3, align=PP_ALIGN.CENTER)

add_speaker_notes(s, """DECISION TREE — make this stick. ~40 sec.

"If you take ONE thing home, this is it. Don't pick a policy from a paper. Pick it from your trace.

Four branches.

If your raw alpha is below 0.5 and your sizes look uniform — Congress-like — pick SIEVE. It's statistically tied with W-TinyLFU and it's simpler. No CMS. No admission filter. Lower implementation cost, lower operational debugging cost, identical performance.

If your raw alpha is above 0.8 with a long tail — Court-like — pick W-TinyLFU. The 5pp gap is stable across every alpha. The CMS earns its keep.

If you have ANY size outlier above 100x median — even if alpha is moderate — pick W-TinyLFU AND report byte miss-ratio with 5-seed CI bands. Otherwise the byte numbers will lie to you.

If your cache is below 1% of working set — admission quality dominates regardless. W-TinyLFU beats LRU by 9 to 13 points on both workloads in that regime.

The policy choice is five lines of code. The hard work is reading the trace right."

(~40 sec)""")

# ===========================================================
# SLIDE 12 — CLOSE
# ===========================================================
s = prs.slides.add_slide(blank); set_bg(s)
page_chrome(s, 12, TOTAL, "ONE SENTENCE  ·  RELATED WORK  ·  THANKS")

# big sentence
add_text(s, Inches(0.5), Inches(1.10), Inches(12.33), Inches(0.4),
         "THE FINDING IN ONE SENTENCE",
         size=11, bold=True, color=ACCENT, letter_space=300)

add_runs(s, Inches(0.5), Inches(1.55), Inches(12.33), Inches(2.8), [
    [("Choice of admission mechanism matters ",
      {"size": 38, "bold": True, "color": INK})],
    [("only ", {"size": 38, "italic": True, "bold": True, "color": ACCENT}),
     ("when the raw workload provides", {"size": 38, "bold": True, "color": INK})],
    [("a frequency gradient to exploit.",
      {"size": 38, "bold": True, "color": INK})],
], line_spacing=1.10)

# horizontal rule
add_line(s, Inches(0.5), Inches(4.40), Inches(12.83), Inches(4.40),
         color=RULE, weight=Pt(0.75))

# Related work, in two columns
add_text(s, Inches(0.5), Inches(4.55), Inches(6), Inches(0.4),
         "WHERE THIS SITS IN THE LITERATURE",
         size=11, bold=True, color=ACCENT, letter_space=250)

add_runs(s, Inches(0.5), Inches(4.95), Inches(6.0), Inches(1.95), [
    [("Yang  S3-FIFO  SOSP '23", {"size": 12, "color": INK, "font": MONO}),
     ("    →  ablated 5/10/20%", {"size": 11, "color": INK_3})],
    [("Zhang  SIEVE  NSDI '24", {"size": 12, "color": INK, "font": MONO}),
     ("      →  visited-bit confirmed +15.4 pp", {"size": 11, "color": INK_3})],
    [("Einziger  TinyLFU  TOS '17", {"size": 12, "color": INK, "font": MONO}),
     (" →  Doorkeeper effect = noise", {"size": 11, "color": INK_3})],
    [("Caffeine  production  '24", {"size": 12, "color": INK, "font": MONO}),
     ("    →  W=1% / SLRU=99% replicated", {"size": 11, "color": INK_3})],
    [("Waldspurger  SHARDS  FAST '15", {"size": 12, "color": INK, "font": MONO}),
     (" →  self-conv. MAE 0.038 < 0.05", {"size": 11, "color": INK_3})],
    [("Clauset  α MLE  SIAM '09", {"size": 12, "color": INK, "font": MONO}),
     ("    →  used for raw α fits", {"size": 11, "color": INK_3})],
])

# right: contributions
add_text(s, Inches(7.00), Inches(4.55), Inches(6), Inches(0.4),
         "WHAT THIS WORK ADDS",
         size=11, bold=True, color=ACCENT, letter_space=250)

add_runs(s, Inches(7.00), Inches(4.95), Inches(5.83), Inches(1.95), [
    [("◆", {"size": 14, "color": ACCENT}),
     ("  Two real legislative-API traces, characterized side-by-side.",
      {"size": 13, "color": INK_2})],
    [("◆", {"size": 14, "color": ACCENT}),
     ("  Welch's-t-backed 5-seed comparison with ±1σ bands.",
      {"size": 13, "color": INK_2})],
    [("◆", {"size": 14, "color": ACCENT}),
     ("  Mechanism explanation grounded in raw α × size dist.",
      {"size": 13, "color": INK_2})],
    [("◆", {"size": 14, "color": ACCENT}),
     ("  Decision tree practitioners can apply in 2 minutes.",
      {"size": 13, "color": INK_2})],
])

# closing footer line
add_line(s, Inches(0.5), Inches(7.00), Inches(12.83), Inches(7.00),
         color=RULE, weight=Pt(0.5))
add_text(s, Inches(0.5), Inches(7.10), Inches(8), Inches(0.30),
         "Mira Yu  ·  github.com/civicache  ·  CS 2640 Spring 2026",
         size=11, color=INK_2, font=MONO, letter_space=120)
add_text(s, Inches(11.0), Inches(7.10), Inches(2.0), Inches(0.30),
         "QUESTIONS",
         size=11, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT,
         letter_space=300)

add_speaker_notes(s, """CLOSE — strong landing. ~30 sec.

"One sentence. Choice of admission mechanism matters ONLY when the raw workload provides a frequency gradient to exploit.

Where this sits — six papers I lean on. S3-FIFO from SOSP'23, ablated. SIEVE from NSDI'24, visited-bit empirically confirmed as load-bearing. TinyLFU from TOS '17, with the Doorkeeper variant tested and shown to be noise. Caffeine's production W=1%, SLRU=99% replicated. SHARDS from FAST'15, MAE 0.038 below the 0.05 sanity gate. Clauset 2009 for the alpha MLE.

Four contributions. Two real traces side-by-side. Welch's-t backed comparison with sigma bands. A mechanism explanation grounded in raw alpha and size distribution. And a 2-minute decision tree.

Questions?"

(~30 sec  — TOTAL TARGET 9:50–10:00)""")

# ---------- save ----------
prs.save(OUT)
print(f"\nWrote {OUT}")
print(f"Slides: {len(prs.slides)}")
